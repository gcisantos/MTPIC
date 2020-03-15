class ObtemTexto:

    #Importação dos pacotes necessários
    import os
    from django.core.files.base import ContentFile
    from datetime import datetime as horaAtual
    from . import models
    from time import sleep

    #Obtem pasta pardrão para referenciar
    config = models.Configuracoes.objects.get(id_config = 1)
    caminho = str(config.caminho_raiz_programas)
    caminho = caminho.replace("\\MTPIC_Processador","")


    #BASE_PATH é a pasta padrão onde se econtra a psta de upload
    BASE_PATH = caminho+'\\MTPIC\\processamento\\'     

    # Variavel que contem parametros passados pelo ajax
    request = None

    def __init__(self, requisicao):
        self.request = requisicao

    def upload_arquivo(self,arquivoEnviado):
        #Define a pasta onde é realizado o upload dos arquivos
        pasta = 'upload_files\\' 

        extensaoArquivo = arquivoEnviado.name.split(".")
        
        
   

        #Realiza o salvamento do arquivo juntando seu nome a e data/hora atual
        nomeCompletoArquivo = self.os.path.join(self.BASE_PATH, pasta, str(self.horaAtual.now().year)+str(self.horaAtual.now().month)+str(self.horaAtual.now().day)+str(self.horaAtual.now().hour)+str(self.horaAtual.now().minute)+str(self.horaAtual.now().second)+str(self.horaAtual.now().microsecond)+"."+extensaoArquivo[-1])
        
        #Cria um arquivo para escrever nele 
        arquivoEscrever = open(nomeCompletoArquivo, 'wb+')

        #obtem conteúdo do arquivo
        file_content = self.ContentFile(arquivoEnviado.read())


        #realzia loop no arquivo escrevendo o mesmo
        print(self.dataHoraAtual()+"Iniciando upload...")
        for chunk in file_content.chunks():            
            arquivoEscrever.write(chunk)                        
        arquivoEscrever.close()
        print(self.dataHoraAtual()+"Upload finalizado...")
        return nomeCompletoArquivo


    def dataHoraAtual(self):
        return "["+str(self.horaAtual.now().hour)+":"+str(self.horaAtual.now().minute)+":"+str(self.horaAtual.now().second)+":"+str(self.horaAtual.now().microsecond)+"] "
    

    def obterTexto(self):
        if self.request.method == 'POST':
            texto = ""
            link = ""
            taxaAceitacao = self.request.POST['taxaAceitacao']
            euristicaProcessamento = self.request.POST['euristicaProcessamento']
            formaBusca = self.request.POST['formaBusca']
            resultado_upload =""
            if formaBusca == "Link":
                #obtem texto do link solicitado e armazena na variavel texto
                link = self.request.POST['link']
                texto = self.obtemTextoLink(link)

            elif formaBusca =="Arquivo":
                # Processa arquivo e extrai texto

                resultado_upload = self.upload_arquivo(self.request.FILES['arquivo'])
                if resultado_upload != "":
                    texto = self.verificaProcessoArquivo(resultado_upload)

            elif formaBusca =="Texto":
                texto = self.request.POST['texto']
            
            linguagemSeleciona =""
            with open("C:\\MTPIC\\MTPIC\\processamento\\static\\config\\config.txt", 'r',encoding='utf8') as configuracoesLeitura:
                for linha in configuracoesLeitura.readlines():
                    if(linha.startswith("Linguagem")):
                        linguagemSeleciona = linha.split(":")
                        linguagemSeleciona = linguagemSeleciona[1]
                configuracoesLeitura.close()


            if(texto!=""):
                busca = self.models.Buscas(
                    forma_busca = formaBusca,
                    taxaAceitacao= int(taxaAceitacao),
                    euristicaProcessamento = int(euristicaProcessamento),
                    caminho_arquivo = resultado_upload,
                    caminho_link = link,
                    tipoBusca = 1,
                    idioma = linguagemSeleciona,
                    texto_busca =texto,
                    dataIni_busca = self.horaAtual.now()
                    )
            else:
                busca = self.models.Buscas(
                    forma_busca = formaBusca,
                    taxaAceitacao= int(taxaAceitacao),
                    euristicaProcessamento = int(euristicaProcessamento),
                    caminho_arquivo = resultado_upload,
                    caminho_link = link,
                    tipoBusca = 1,
                    idioma = linguagemSeleciona,
                    texto_busca =texto,
                    processado_busca = 1,
                    dataIni_busca = self.horaAtual.now()
                    )
                print("Não foi possivel extratir texto")
                return "0"
            busca.save()

            id_processado = busca.id_busca
            
            processou = True
            while(processou):
                self.sleep(0.05)                
                buscaProcessada = self.models.Buscas.objects.get(id_busca = id_processado)
                if(buscaProcessada.processado_busca > 0):                   
                    return buscaProcessada.resultadoTema_busca
        return ""
       

       
    def obtemSubTema(self):
        if self.request.method == 'POST':
            texto = ""
            link = ""
            taxaAceitacao =85
            euristicaProcessamento = 1
            formaBusca = self.request.POST['formaBusca']
            resultado_upload =""
            if formaBusca == "Link":
                #obtem texto do link solicitado e armazena na variavel texto
                link = self.request.POST['link']
                texto = self.obtemTextoLink(link)

            elif formaBusca =="Arquivo":
                # Processa arquivo e extrai texto

                resultado_upload = self.upload_arquivo(self.request.FILES['arquivo'])
                if resultado_upload != "":
                    texto = self.verificaProcessoArquivo(resultado_upload)

            elif formaBusca =="Texto":
                texto = self.request.POST['texto']
            
            linguagemSeleciona =""
            with open("C:\\MTPIC\\MTPIC\\processamento\\static\\config\\config.txt", 'r',encoding='utf8') as configuracoesLeitura:
                for linha in configuracoesLeitura.readlines():
                    if(linha.startswith("Linguagem")):
                        linguagemSeleciona = linha.split(":")
                        linguagemSeleciona = linguagemSeleciona[1]
                configuracoesLeitura.close()

            if(texto!=""):   
                busca = self.models.Buscas(
                    forma_busca = formaBusca,
                    taxaAceitacao= int(taxaAceitacao),
                    euristicaProcessamento = int(euristicaProcessamento),
                    caminho_arquivo = resultado_upload,
                    caminho_link = link,
                    tipoBusca = 2,
                    idioma = linguagemSeleciona,
                    texto_busca =texto,
                    dataIni_busca = self.horaAtual.now()
                    )
            else:
                busca = self.models.Buscas(
                forma_busca = formaBusca,
                taxaAceitacao= int(taxaAceitacao),
                euristicaProcessamento = int(euristicaProcessamento),
                caminho_arquivo = resultado_upload,
                caminho_link = link,
                tipoBusca = 2,
                idioma = linguagemSeleciona,
                processado_busca = 1,
                texto_busca =texto,
                dataIni_busca = self.horaAtual.now()
                )
                print("Não foi possivel extratir texto")
                return "0"

            busca.save()

            id_processado = busca.id_busca
            
            processou = True
            while(processou):
                self.sleep(0.05)                
                buscaProcessada = self.models.Buscas.objects.get(id_busca = id_processado)
                if(buscaProcessada.processado_busca > 0):    
                    processou = False               
                    return buscaProcessada.resultadoAssuntos_proximidades_busca
        return ""



    def obtemSubTema_Tema(self):
        if self.request.method == 'POST':
            texto = ""
            link = ""
            taxaAceitacao = self.request.POST['taxaAceitacao']
            euristicaProcessamento = self.request.POST['euristicaProcessamento']
            formaBusca = self.request.POST['formaBusca']
            resultado_upload =""
            if formaBusca == "Link":
                #obtem texto do link solicitado e armazena na variavel texto
                link = self.request.POST['link']
                texto = self.obtemTextoLink(link)

            elif formaBusca =="Arquivo":
                # Processa arquivo e extrai texto

                resultado_upload = self.upload_arquivo(self.request.FILES['arquivo'])
                if resultado_upload != "":
                    texto = self.verificaProcessoArquivo(resultado_upload)

            elif formaBusca =="Texto":
                texto = self.request.POST['texto']
            
            linguagemSeleciona =""
            with open("C:\\MTPIC\\MTPIC\\processamento\\static\\config\\config.txt", 'r',encoding='utf8') as configuracoesLeitura:
                for linha in configuracoesLeitura.readlines():
                    if(linha.startswith("Linguagem")):
                        linguagemSeleciona = linha.split(":")
                        linguagemSeleciona = linguagemSeleciona[1]
                configuracoesLeitura.close()

            if(texto!=""): 
                busca = self.models.Buscas(
                forma_busca = formaBusca,
                taxaAceitacao= int(taxaAceitacao),
                euristicaProcessamento = int(euristicaProcessamento),
                caminho_arquivo = resultado_upload,
                caminho_link = link,
                tipoBusca = 3,
                idioma = linguagemSeleciona,
                texto_busca =texto,
                dataIni_busca = self.horaAtual.now()
                )
            else:
                busca = self.models.Buscas(
                forma_busca = formaBusca,
                taxaAceitacao= int(taxaAceitacao),
                euristicaProcessamento = int(euristicaProcessamento),
                caminho_arquivo = resultado_upload,
                caminho_link = link,
                tipoBusca = 3,
                idioma = linguagemSeleciona,
                processado_busca=1,
                texto_busca =texto,
                dataIni_busca = self.horaAtual.now()
                )
                print("Não foi possivel extratir texto")
                return "","","","1"



            busca.save()

            id_processado = busca.id_busca
            
            processou = True
            while(processou):
                self.sleep(1)                
                buscaProcessada = self.models.Buscas.objects.get(id_busca = id_processado)
                if(buscaProcessada.processado_busca > 0):                                     
                    return (buscaProcessada.resultadoTema_busca,buscaProcessada.resultadoAssuntos_proximidades_busca,buscaProcessada.resultado_Grafo,"0")
        return ""
       
    def obtemSubTema_Tema_processados(self):
        if self.request.method == 'POST':
            id_processado = self.request.POST['idBusca']
            buscaProcessada = self.models.Buscas.objects.filter(id_busca = id_processado)
            if(buscaProcessada.count() > 0):
                result = buscaProcessada[0]
                if(result.processado_busca > 0):                                     
                    return (result.resultadoTema_busca,result.resultadoAssuntos_proximidades_busca,result.resultado_Grafo,"0")
        return "","","","1"

    def obtemTextoLink(self, link):
        #Realiza importação das bibliotecas necessárias
        from bs4 import BeautifulSoup
        import requests
        
        #realiza requisição do site
        response = requests.get(link)
        
        #Realiza o parser do html
        soup = BeautifulSoup(response.content, "html.parser")

        #Verefica se o resultado obtido não é nulo
        if(soup.body  is not None):
            #obtem os textos pertencentes as tags 'p'
            textoExtraido  = [texto.get_text().replace('\n','') for texto in soup.body.find_all('p') ]
            
            #retorna o texto do site junto ao seu titulo
            return str(soup.title.get_text()) + "\n".join(textoExtraido)
        #Caso o resultado do site seja nulo retorna vazio
        return ""
    

    def verificaProcessoArquivo(self,caminho):
        import os
        name, ext = os.path.splitext(caminho)
        extensao = ext.replace(".","")
        
        #Processa arquivo word
        if extensao == "docx" or extensao=="dotx":
            return self.obtemTextoDocx(caminho)
        #processa imagem
        elif(extensao=="jpg" or extensao=="jpeg" or extensao=="png"or extensao=="jpe"or extensao=="bmp"):
            return self.obtemTextoImagem(caminho)
        #processa PDF
        elif(extensao =="pdf"):
            return self.obtemTextoPDF(caminho)
        #processa pptx
        elif(extensao =="pptx"):
            return self.obtemTextoPPTX(caminho)
        #processa txt
        elif(extensao == "txt"):
            return self.obtemTextoTXT(caminho)
        return ""


    def obtemTextoDocx(self, caminhoDocumento):
        import docxpy        
        return docxpy.process(caminhoDocumento)

    def obtemTextoImagem(self, caminhoImagem):
        import pytesseract as ocr
        from PIL import Image
        ocr.pytesseract.tesseract_cmd = "C:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe"
        return ocr.image_to_string(Image.open(caminhoImagem), lang='por')

    def obtemTextoPDF(self, caminhoPDF):
        import fitz
        doc = fitz.open(caminhoPDF)
        textos = [page.getText() for page in doc]
        retorno = ' '.join(textos)
        if(retorno!=""):
            return retorno
        cont = 0
        for page in doc:
            pix = page.getPixmap()
            pix.writePNG(self.BASE_PATH+"upload_files\\page-"+str(cont)+".png")
            retorno = retorno+self.obtemTextoImagem(self.BASE_PATH+"upload_files\\page-"+str(cont)+".png")
        return retorno

    def obtemTextoPPTX(self, caminhoPPTX):
        from pptx import Presentation
        textos = []
        ppt = Presentation(caminhoPPTX)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    textos.append(shape.text)
        return ' '.join(textos)

    def obtemTextoTXT(self, caminhoTXT):
        arquivo = open(caminhoTXT, "r") 
        texto = arquivo.read()
        arquivo.close()
        return texto
